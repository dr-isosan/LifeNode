#!/usr/bin/env python3
"""
LifeNode Sunum Oluşturucu
Python-pptx ile PowerPoint sunumu oluşturur
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle, author=""):
    """Kapak slaydı ekle"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Arka plan rengi
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(44, 62, 80)  # Koyu mavi
    background.line.fill.background()

    # Başlık
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Alt başlık
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(9), Inches(0.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(230, 126, 34)  # Turuncu
    p.alignment = PP_ALIGN.CENTER

    # Yazar
    if author:
        author_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(9), Inches(0.5)
        )
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = author
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(189, 195, 199)
        p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content_items, is_bullet=True):
    """İçerik slaydı ekle"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Başlık şeridi
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(44, 62, 80)
    header.line.fill.background()

    # Başlık
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # İçerik
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if is_bullet and item.strip():
            p.text = f"• {item}"
        else:
            p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_after = Pt(12)


def add_two_column_slide(prs, title, left_items, right_items):
    """İki sütunlu slayt ekle"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Başlık şeridi
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(44, 62, 80)
    header.line.fill.background()

    # Başlık
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Sol sütun
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4.3), Inches(5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_after = Pt(8)

    # Sağ sütun
    right_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.5), Inches(4.3), Inches(5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_after = Pt(8)


def add_table_slide(prs, title, headers, rows):
    """Tablo içeren slayt ekle"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Başlık şeridi
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(44, 62, 80)
    header_shape.line.fill.background()

    # Başlık
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Tablo
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * num_rows)
    ).table

    # Başlık satırı
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Veri satırları
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(236, 240, 241)


def add_thank_you_slide(prs, title="Teşekkür Ederiz", subtitle="Sorularınız?"):
    """Teşekkür slaydı ekle"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Arka plan
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(44, 62, 80)
    background.line.fill.background()

    # Başlık
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Alt başlık
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.8), Inches(9), Inches(0.6)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(230, 126, 34)
    p.alignment = PP_ALIGN.CENTER


def create_lifenode_presentation():
    """LifeNode sunumunu oluştur"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Kapak Slaydı
    add_title_slide(
        prs,
        "LifeNode 🌐",
        "AI-Driven Dynamic Routing Simulation",
        "Samsung Innovation Campus\nİshak Duran",
    )

    # 2. Giriş ve Arka Plan
    add_content_slide(
        prs,
        "Giriş ve Arka Plan",
        [
            "Afet koşullarında (deprem, sel, savaş) merkezi iletişim altyapıları çökebilir.",
            "GSM kuleleri, fiber optik hatlar ve internet altyapısı tek noktadan arızaya açıktır.",
            "Kurtarma ekipleri ve kazazedeler arasındaki iletişim hayati önem taşır.",
            "Bu proje, merkezi olmayan ad-hoc mesh ağlarını simüle ederek yapay zeka tabanlı dinamik yönlendirme çözümleri geliştirmeyi amaçlamaktadır.",
        ],
    )

    # 3. Projenin Hedefleri
    add_content_slide(
        prs,
        "Projenin Hedefleri",
        [
            "Dinamik Yönlendirme – Hareket eden düğümler arasında kesintisiz iletişim sağlama.",
            "Yapay Zeka Entegrasyonu – Q-Learning tabanlı RL agent ile değişen koşullara adaptasyon.",
            "Karşılaştırmalı Analiz – Dijkstra, AODV ve Q-Learning algoritmalarını karşılaştırma.",
            "Afet Simülasyonu – Deprem, kademeli arıza senaryolarında ağ dayanıklılığını test etme.",
            "Metrik Analizi – PDR, Latency, Hop Count gibi kritik performans ölçümlerini raporlama.",
        ],
    )

    # 4. SDG Hedefleri
    add_content_slide(
        prs,
        "Sürdürülebilir Kalkınma Hedefleri (SDG)",
        [
            "SDG 9: Endüstri, Yenilik ve Altyapı – Dayanıklı altyapı ve yenilikçi AI uygulamaları.",
            "SDG 11: Sürdürülebilir Şehirler – Afetlere karşı dirençli iletişim sistemleri.",
            "SDG 13: İklim Eylemi – Afet müdahale kapasitesini artırma.",
            "SDG 17: Ortaklıklar – Açık kaynak araştırma projesi olarak bilgi paylaşımı.",
        ],
    )

    # 5. Kullanılan Teknolojiler
    add_content_slide(
        prs,
        "Kullanılan Teknolojiler ve Modeller",
        [
            "Simülasyon Ortamı – Python tabanlı modüler mimari ile düğüm, bağlantı ve paket simülasyonu.",
            "Yönlendirme Algoritmaları – Dijkstra (baseline), AODV (reaktif) ve Q-Learning tabanlı RL Router.",
            "AODV Protokolü – Ad-hoc On-Demand Distance Vector, ihtiyaç duyulduğunda rota keşfi.",
            "Reinforcement Learning – Tabular Q-Learning algoritması ile dinamik karar verme.",
            "Metrik Toplama – PDR, Latency, Hop Count analizi ve görselleştirme.",
            "Afet Senaryoları – Deprem, kademeli arıza simülasyonları.",
        ],
    )

    # 6. Metodoloji
    add_two_column_slide(
        prs,
        "Metodoloji",
        [
            "Veri Toplama ve İşleme:",
            "  - 40 düğümlü ağ topolojisi",
            "  - 500 adımlık simülasyon",
            "",
            "Model Eğitimi:",
            "  - 30 episode Q-Learning",
            "  - Learning rate: 0.1",
            "  - Discount factor: 0.9",
        ],
        [
            "Gerçek Zamanlı Analiz:",
            "  - Paket yönlendirme performansı",
            "  - Enerji tüketimi takibi",
            "",
            "Karşılaştırmalı Test:",
            "  - Dijkstra vs AODV vs RL",
            "  - Normal ve afet senaryoları",
        ],
    )

    # 7. RL Agent Mimarisi
    add_two_column_slide(
        prs,
        "RL Agent Mimarisi",
        [
            "State (Durum):",
            "  - Düğümün enerji seviyesi",
            "  - Kuyruk doluluğu",
            "  - Komşu sinyal kalitesi",
            "",
            "Action (Eylem):",
            "  - Paketi hangi komşuya ileteceğini seçme",
        ],
        [
            "Reward (Ödül):",
            "  ✅ Başarılı teslimat: +10 puan",
            "  ❌ Paket kaybı: -10 puan",
            "  🐢 Gecikme/Hop cezası: -0.1 puan",
            "",
            "Algoritma:",
            "  - Tabular Q-Learning",
            "  - Epsilon-greedy exploration",
        ],
    )

    # 8. Çıktılar ve Sonuçlar
    add_table_slide(
        prs,
        "Çıktılar ve Sonuçlar",
        ["Senaryo", "Router", "PDR", "Latency", "Hop"],
        [
            ["Normal", "Dijkstra", "%12.6", "7.3ms", "2.1"],
            ["Normal", "AODV", "%14.1", "7.8ms", "2.4"],
            ["Normal", "RL Agent", "%15.2", "6.8ms", "2.3"],
            ["Deprem", "Dijkstra", "%8.4", "9.1ms", "2.5"],
            ["Deprem", "AODV", "%10.2", "8.9ms", "2.7"],
            ["Deprem", "RL Agent", "%11.7", "8.2ms", "2.4"],
            ["Şiddetli", "Dijkstra", "%5.2", "11.3ms", "2.8"],
            ["Şiddetli", "AODV", "%7.1", "10.5ms", "3.0"],
            ["Şiddetli", "RL Agent", "%8.9", "9.7ms", "2.6"],
        ],
    )

    # 9. Metrikler
    add_table_slide(
        prs,
        "Metrikler ve Başarı Kriterleri",
        ["Metrik", "Açıklama", "Hedef"],
        [
            ["PDR", "Packet Delivery Ratio", "%100'e yakın"],
            ["Latency", "Uçtan uca gecikme (ms)", "Minimum"],
            ["Hop Count", "Paketin geçtiği düğüm sayısı", "Optimum yol"],
            ["Recovery Time", "Arıza sonrası iyileşme", "En hızlı"],
        ],
    )

    # 10. Kullanım Senaryoları
    add_two_column_slide(
        prs,
        "Kullanım Senaryoları",
        [
            "Afet Müdahale Ekipleri:",
            "  Deprem, sel gibi afetlerde",
            "  kurtarma ekipleri arası iletişim",
            "",
            "Askeri Operasyonlar:",
            "  Altyapısız bölgelerde",
            "  güvenli mesh ağ iletişimi",
        ],
        [
            "Uzak Bölgeler:",
            "  İnternet altyapısı olmayan",
            "  kırsal alanlarda iletişim",
            "",
            "IoT Ağları:",
            "  Sensör ağlarında dinamik",
            "  yönlendirme optimizasyonu",
        ],
    )

    # 11. Proje Mimarisi
    add_content_slide(
        prs,
        "Proje Mimarisi",
        [
            "lifenode/",
            "├── environment/     # Simülasyon ortamı (Node, Link, Packet)",
            "├── routing/         # Yönlendirme protokolleri (Dijkstra, AODV)",
            "├── rl_agent/        # Reinforcement Learning (Q-Learning)",
            "├── metrics/         # Performans ölçümü ve görselleştirme",
            "├── scenarios/       # Afet senaryoları (Deprem, Kademeli Arıza)",
            "└── experiments/     # Deney çalıştırma ve raporlama",
        ],
        is_bullet=False,
    )

    # 12. Sonraki Adımlar
    add_content_slide(
        prs,
        "Sonraki Adımlar ve Geliştirme Alanları",
        [
            "Deep Q-Network (DQN) – Derin öğrenme ile daha karmaşık durumları yönetme.",
            "Gerçek Donanım Testi – ESP32 veya LoRa modülleriyle fiziksel mesh ağ.",
            "Enerji Hasadı – Solar panelli düğümlerin simülasyonu.",
            "Büyük Ölçekli Test – 100+ düğümlü ağlarda performans analizi.",
            "Hibrit Yönlendirme – Dijkstra + RL kombinasyonu ile optimal sonuçlar.",
        ],
    )

    # 13. Sonuç
    add_content_slide(
        prs,
        "Sonuç",
        [
            "Üç algoritma karşılaştırıldı: Dijkstra (proaktif), AODV (reaktif), RL (adaptif).",
            "Sonuçlar: RL Agent > AODV > Dijkstra sıralamasıyla performans gösterdi.",
            "AODV, reaktif yapısıyla Dijkstra'dan daha iyi adapte oldu.",
            "RL Agent, özellikle afet senaryolarında en iyi performansı sergiledi.",
            "Sistem modüler yapısıyla kolayca genişletilebilir.",
        ],
    )

    # 14. Teşekkür Slaydı
    add_thank_you_slide(
        prs,
        "Teşekkür Ederiz 👏",
        "LifeNode 🌐 - Dayanıklı İletişim İçin Akıllı Çözümler\n\nSorularınız?",
    )

    # Kaydet
    output_path = "LifeNode_Sunum.pptx"
    prs.save(output_path)
    print(f"✅ Sunum başarıyla oluşturuldu: {output_path}")
    return output_path


if __name__ == "__main__":
    create_lifenode_presentation()
